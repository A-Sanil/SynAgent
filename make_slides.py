"""
Build SynAgent_Aaditya_Presentation.pptx using python-pptx.
Run: python make_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

# ── Helpers ──────────────────────────────────────────────────────────────────

def rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def bg(slide, color_hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)

def box(slide, x, y, w, h, fill_hex=None, line_hex=None, line_width=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_hex)
    else:
        shape.fill.background()
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=14, bold=False, italic=False, color="1E293B",
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
        font="Calibri", wrap=True, margin=0.05):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    tf.margin_left  = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top   = Inches(margin * 0.5)
    tf.margin_bottom = Inches(margin * 0.5)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return txBox

def txt_lines(slide, lines, x, y, w, h,
              size=13, bold=False, color="1E293B",
              align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
              font="Calibri", line_spacing=None, margin=0.08):
    """lines: list of (text, bold, color, size) or just str"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left   = Inches(margin)
    tf.margin_right  = Inches(margin)
    tf.margin_top    = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = valign

    first = True
    for line in lines:
        if isinstance(line, str):
            ltext, lbold, lcolor, lsize = line, bold, color, size
        else:
            ltext = line[0]
            lbold  = line[1] if len(line) > 1 else bold
            lcolor = line[2] if len(line) > 2 else color
            lsize  = line[3] if len(line) > 3 else size

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            from pptx.util import Pt as _Pt
            p.line_spacing = _Pt(line_spacing)
        run = p.add_run()
        run.text = ltext
        run.font.name = font
        run.font.size = Pt(lsize)
        run.font.bold = lbold
        run.font.color.rgb = rgb(lcolor)
    return txBox


# ── Palette ───────────────────────────────────────────────────────────────────
NAVY     = "1F3864"
TEAL     = "0D9488"
TEAL_LT  = "CCFBF1"
WHITE    = "FFFFFF"
LIGHT_BG = "F0F4F8"
CARD_BG  = "FFFFFF"
DARK     = "1E293B"
MID      = "475569"
MUTED    = "94A3B8"
BLUE_BOX = "1E3A5F"
TEAL_BOX = "0A5C5C"
PURP_BOX = "3B1F6B"
ACCENT   = "0D9488"
GREEN    = "059669"
RED_MUT  = "DC2626"
GOLD     = "D97706"

# ── Presentation setup ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = 13.33
H = 7.5
blank_layout = prs.slide_layouts[6]  # completely blank

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, NAVY)

# Teal accent bar on left edge
box(s, 0, 0, 0.18, H, TEAL)

# Decorative teal circle top-right (subtle)
circ = s.shapes.add_shape(9, Inches(11.5), Inches(-1.2), Inches(3.5), Inches(3.5))  # OVAL
circ.fill.solid(); circ.fill.fore_color.rgb = rgb("0F4C5C")
circ.line.fill.background()

circ2 = s.shapes.add_shape(9, Inches(10.8), Inches(5.5), Inches(2.5), Inches(2.5))
circ2.fill.solid(); circ2.fill.fore_color.rgb = rgb("0D3D4F")
circ2.line.fill.background()

# Main title
txt(s, "SynAgent:", 0.5, 1.8, 12, 1.1,
    size=54, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.CENTER)
txt(s, "Patent Data Pipeline", 0.5, 2.8, 12, 1.0,
    size=40, bold=False, color="7DD3FC", font="Georgia", align=PP_ALIGN.CENTER)

# Subtitle
txt(s, "Mining Chemical Synthesis Data from Unstructured Patent Text",
    0.8, 4.05, 11.5, 0.65,
    size=16, bold=False, color="CBD5E1", font="Calibri", align=PP_ALIGN.CENTER)

# Bottom line
box(s, 0.5, 5.0, 12.3, 0.02, TEAL)

# Contact
txt(s, "Aaditya N Sanil  ·  asanil@berkeley.edu  ·  May 2026",
    0.5, 5.2, 12.3, 0.5,
    size=13, color="94A3B8", font="Calibri", align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Gap
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, LIGHT_BG)

# Title area
box(s, 0, 0, W, 1.05, NAVY)
txt(s, "The Gap in Retrosynthesis Tools", 0.5, 0.15, 12.3, 0.75,
    size=30, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.LEFT)

# LEFT CARD — What SynAgent does great
box(s, 0.4, 1.25, 5.8, 3.8, "E8F5F3", "0D9488", 1)
txt(s, "✓  What SynAgent Does Great", 0.55, 1.35, 5.5, 0.5,
    size=14, bold=True, color=TEAL, font="Calibri")
txt_lines(s, [
    ("Validates reaction SMARTS and SMILES via RDKit", False, DARK, 13),
    ("Prices building blocks via ChemSpace API", False, DARK, 13),
    ("Scores route hazard using PubChem GHS codes", False, DARK, 13),
    ("Gemini LLM orchestrates 4 parallel specialist agents", False, DARK, 13),
    ("Handles full retrosynthesis route evaluation", False, DARK, 13),
], 0.65, 1.95, 5.4, 2.9, font="Calibri", line_spacing=22, margin=0.05)

# RIGHT CARD — What was missing
box(s, 6.9, 1.25, 5.8, 3.8, "FEF2F2", "DC2626", 1)
txt(s, "✗  What Was Missing", 7.05, 1.35, 5.5, 0.5,
    size=14, bold=True, color=RED_MUT, font="Calibri")
txt_lines(s, [
    ("No experimental precedent for proposed reactions", False, DARK, 13),
    ("No yield data from real literature", False, DARK, 13),
    ("No connection to patent databases", False, DARK, 13),
    ("Routes scored on price & hazard — not feasibility", False, DARK, 13),
    ("No way to ask: has this reaction been done before?", False, DARK, 13),
], 7.1, 1.95, 5.5, 2.9, font="Calibri", line_spacing=22, margin=0.05)

# Bottom callout box
box(s, 0.4, 5.3, 12.5, 1.0, "FFF7ED", GOLD, 1)
txt(s, "Goal: Add a real reaction corpus so SynAgent can answer — has this reaction been done before, and what yield did it achieve?",
    0.6, 5.4, 12.1, 0.8,
    size=14, bold=False, color="92400E", font="Calibri", align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Three Systems
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, LIGHT_BG)

box(s, 0, 0, W, 1.05, NAVY)
txt(s, "What I Built — Three Systems", 0.5, 0.15, 12.3, 0.75,
    size=30, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.LEFT)

cards = [
    (BLUE_BOX, "1E3A5F", "📄", "Patent Ingestion\nPipeline",
     "Mine synthesis reactions from chemistry paper abstracts using Gemini LLM.\n\nExtracts: reaction SMARTS, yield, temperature, solvent, catalyst.\n\nRate-limited worker with supervisor loop."),
    (TEAL_BOX, "0A5C5C", "🗄️", "Reaction Database",
     "SQLite + FTS5 full-text search.\n16,831 reactions indexed.\n\nTwo sources:\n• ORD expert-curated (confidence=1.0)\n• LLM-extracted from patents (0–1)"),
    (PURP_BOX, "3B1F6B", "🔍", "Search Frontend",
     "Standalone FastAPI web app on port 8001.\n\nSearch by keyword, SMILES, or SMARTS.\nFilter by yield, source, confidence.\nCard results with one-click copy."),
]

bx = 0.4
for (fill, hdr, icon, title, desc) in cards:
    box(s, bx, 1.25, 3.95, 5.6, fill)
    # Icon circle
    circ = s.shapes.add_shape(9, Inches(bx+1.6), Inches(1.45), Inches(0.75), Inches(0.75))
    circ.fill.solid(); circ.fill.fore_color.rgb = rgb(hdr)
    circ.line.fill.background()
    txt(s, icon, bx+1.6, 1.45, 0.75, 0.75, size=22, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, color=WHITE)
    txt(s, title, bx+0.15, 2.35, 3.65, 0.9,
        size=18, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.CENTER)
    txt(s, desc, bx+0.2, 3.3, 3.55, 3.2,
        size=12.5, color="CBD5E1", font="Calibri", align=PP_ALIGN.LEFT)
    bx += 4.45

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Two Ingestion Paths
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, LIGHT_BG)

box(s, 0, 0, W, 1.05, NAVY)
txt(s, "Two Ways Reactions Enter the Database", 0.5, 0.15, 12.3, 0.75,
    size=30, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.LEFT)

def flow_node(slide, label, x, y, w=1.8, h=0.65, fill="1E3A5F", tcolor=WHITE, tsize=11):
    box(slide, x, y, w, h, fill)
    txt(slide, label, x, y, w, h, size=tsize, color=tcolor, font="Calibri",
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.04)

def arrow_h(slide, x, y, length=0.4):
    line = slide.shapes.add_shape(1, Inches(x), Inches(y+0.28), Inches(length), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = rgb("64748B")
    line.line.fill.background()
    # arrowhead (small triangle)
    arr = slide.shapes.add_shape(5,  # RIGHT_TRIANGLE -- use a thin box instead
        Inches(x+length-0.01), Inches(y+0.2), Inches(0.12), Inches(0.2))
    arr.fill.solid(); arr.fill.fore_color.rgb = rgb("64748B")
    arr.line.fill.background()

# PATH A
box(s, 0.3, 1.2, 12.7, 2.1, "EBF5FF", "1E3A5F", 1)
txt(s, "PATH A — Patent Mining (LLM Extraction)", 0.45, 1.3, 8, 0.4,
    size=12, bold=True, color=BLUE_BOX, font="Calibri")

nodes_a = ["Europe PMC\nAPI", "Fetch\nAbstracts", "parse_queue", "Gemini 2.5\nFlash", "JSON\nExtraction", "reactions\ntable"]
fills_a = ["334155","1E3A5F","1E3A5F","0A5C5C","0A5C5C","065F46"]
xa = 0.45
for i, (n, f) in enumerate(zip(nodes_a, fills_a)):
    flow_node(s, n, xa, 1.78, 1.75, 0.72, f)
    if i < len(nodes_a)-1:
        arrow_h(s, xa+1.75, 1.78)
    xa += 2.15

txt(s, "source='patent' · confidence=0.0–1.0", 0.5, 2.6, 12.2, 0.35,
    size=11, color=MID, font="Calibri", italic=True)

# PATH B
box(s, 0.3, 3.55, 12.7, 2.0, "F0FDF4", "059669", 1)
txt(s, "PATH B — ORD Direct  (No LLM — Expert Curated)", 0.45, 3.65, 9, 0.4,
    size=12, bold=True, color="065F46", font="Calibri")

nodes_b = ["HuggingFace\nORD Dataset", "ord_ingestion\n.py", "reactions\ntable"]
fills_b = ["064E3B","065F46","065F46"]
xa = 0.45
for i, (n, f) in enumerate(zip(nodes_b, fills_b)):
    flow_node(s, n, xa, 4.1, 1.75, 0.72, f)
    if i < len(nodes_b)-1:
        arrow_h(s, xa+1.75, 4.1)
    xa += 2.15

box(s, 4.7, 4.0, 7.8, 0.9, "DCFCE7", "059669", 1)
txt_lines(s, [
    ("No LLM needed — data is pre-structured by expert chemists", True, "065F46", 12),
    ("source='ord' · confidence=1.0 · zero hallucination risk", False, "065F46", 11),
], 4.85, 4.05, 7.5, 0.8, font="Calibri", line_spacing=20, margin=0.05)

txt(s, "↑ Both paths write to the same SQLite reactions table with FTS5 full-text search index",
    0.4, 5.7, 12.5, 0.45,
    size=12, bold=False, color=MID, font="Calibri", align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Gemini Extraction
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, LIGHT_BG)

box(s, 0, 0, W, 1.05, NAVY)
txt(s, "How Gemini Extracts Reactions from Text", 0.5, 0.15, 12.3, 0.75,
    size=30, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.LEFT)

steps = [
    ("1", "Raw Paper Text", "Abstract + methods section fetched from Europe PMC API"),
    ("2", "Prompt Sent to Gemini 2.5 Flash", "\"Extract all chemical reactions as JSON: reaction_smarts, reactants, product, yield_percent, temperature_celsius, solvent, catalyst, confidence\""),
    ("3", "Structured JSON Response", "Gemini returns machine-readable reaction data with confidence score (0–1)"),
    ("4", "RDKit Validation", "SMILES canonicalized, reaction SMARTS validated, invalid structures filtered"),
    ("5", "Stored in Database", "Inserted into reactions table with source='patent', ready for FTS5 search"),
]

y = 1.2
for (num, title, desc) in steps:
    # Number circle
    circ = s.shapes.add_shape(9, Inches(0.35), Inches(y+0.05), Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = rgb(TEAL)
    circ.line.fill.background()
    txt(s, num, 0.35, y+0.05, 0.55, 0.55, size=14, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Georgia")
    # Title
    txt(s, title, 1.05, y, 5.8, 0.38, size=13, bold=True, color=DARK, font="Calibri")
    # Desc
    txt(s, desc, 1.05, y+0.35, 5.8, 0.5, size=11.5, color=MID, font="Calibri")
    # Connector
    if num != "5":
        ln = s.shapes.add_shape(1, Inches(0.58), Inches(y+0.6), Inches(0.08), Inches(0.28))
        ln.fill.solid(); ln.fill.fore_color.rgb = rgb("CBD5E1")
        ln.line.fill.background()
    y += 1.02

# Right side box — rate limit note
box(s, 7.3, 1.2, 5.6, 3.1, "FFF7ED", GOLD, 1)
txt(s, "⚠  Rate Limit Management", 7.5, 1.3, 5.2, 0.45,
    size=13, bold=True, color="92400E", font="Calibri")
txt_lines(s, [
    ("Free tier: 15 RPM / ~1,500 req/day", False, "92400E", 12),
    ("Pacing: 7 RPM (8s heartbeat delay)", False, "92400E", 12),
    ("run_loop.py supervisor auto-restarts", False, "92400E", 12),
    ("on rate-limit failures (429 errors)", False, "92400E", 12),
    ("GEMINI_MODEL env var — swap models", False, "92400E", 12),
    ("without code changes", False, "92400E", 12),
], 7.5, 1.85, 5.2, 2.2, font="Calibri", line_spacing=20, margin=0.05)

box(s, 7.3, 4.5, 5.6, 2.7, "F0FDF4", "059669", 1)
txt(s, "📊  Current Progress", 7.5, 4.6, 5.2, 0.4,
    size=13, bold=True, color="065F46", font="Calibri")
txt_lines(s, [
    ("1,542 papers fetched from Europe PMC", False, "065F46", 12),
    ("1,700+ papers queued for parsing", False, "065F46", 12),
    ("81 reactions extracted (LLM)", False, "065F46", 12),
    ("16,750 reactions from ORD corpus", False, "065F46", 12),
    ("16,831 total reactions in database", True, "065F46", 13),
], 7.5, 5.1, 5.2, 1.95, font="Calibri", line_spacing=20, margin=0.05)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Database Stats
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, LIGHT_BG)

box(s, 0, 0, W, 1.05, NAVY)
txt(s, "Unified Reactions Database", 0.5, 0.15, 12.3, 0.75,
    size=30, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.LEFT)

# Big stat cards
stats = [
    ("16,831", "Total Reactions", TEAL, "E0F2FE"),
    ("16,750", "ORD Expert-Curated", "059669", "DCFCE7"),
    ("5,017",  "With Yield Data", GOLD, "FFF7ED"),
    ("16,770", "High Confidence\n(≥ 0.6)", BLUE_BOX, "EFF6FF"),
]
sx = 0.4
for (num, label, accent, bg_c) in stats:
    box(s, sx, 1.2, 2.95, 1.8, bg_c, accent, 2)
    txt(s, num, sx, 1.35, 2.95, 0.9,
        size=34, bold=True, color=accent, font="Georgia", align=PP_ALIGN.CENTER)
    txt(s, label, sx, 2.25, 2.95, 0.65,
        size=11.5, color=DARK, font="Calibri", align=PP_ALIGN.CENTER)
    sx += 3.15

# Schema table
headers = ["Column", "Type", "Description"]
rows = [
    ["reaction_smarts", "TEXT", "Full reaction template (SMARTS notation)"],
    ["product_smiles",  "TEXT", "Product SMILES string"],
    ["yield_percent",   "REAL", "Experimental yield (%)"],
    ["temperature_celsius", "REAL", "Reaction temperature"],
    ["solvent",         "TEXT", "Solvent SMILES"],
    ["catalyst",        "TEXT", "Catalyst SMILES"],
    ["confidence",      "REAL", "0–1 score (LLM certainty)"],
    ["source",          "TEXT", "'ord' (expert) or 'patent' (LLM)"],
]

# Table header
box(s, 0.4, 3.2, 12.5, 0.42, NAVY)
col_ws = [2.8, 1.2, 8.5]
cx = 0.5
for (h, cw) in zip(headers, col_ws):
    txt(s, h, cx, 3.22, cw, 0.38, size=11, bold=True, color=WHITE, font="Calibri")
    cx += cw

for i, row in enumerate(rows):
    row_bg = "F8FAFC" if i % 2 == 0 else WHITE
    box(s, 0.4, 3.62+i*0.44, 12.5, 0.44, row_bg, "E2E8F0", 0)
    cx = 0.5
    for (cell, cw) in zip(row, col_ws):
        bold_c = TEAL if cx < 1 else DARK
        mono = cx < 3.5
        txt(s, cell, cx, 3.64+i*0.44, cw, 0.4,
            size=11, bold=(i==7 and cx>3), color=DARK,
            font="Courier New" if mono else "Calibri")
        cx += cw

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Search Frontend Mockup
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, LIGHT_BG)

box(s, 0, 0, W, 1.05, NAVY)
txt(s, "Reaction Search UI  —  localhost:8001", 0.5, 0.15, 12.3, 0.75,
    size=30, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.LEFT)

# App mockup frame
box(s, 0.35, 1.15, 12.6, 6.1, "0F1117", "334155", 1)

# App header bar
box(s, 0.35, 1.15, 12.6, 0.52, "0D1120")
txt(s, "⚗️  Reaction Search  ·  SynAgent  ·  Patent & ORD Database",
    0.55, 1.2, 8, 0.42, size=11, bold=True, color=WHITE, font="Calibri")

# Stats strip
box(s, 0.35, 1.67, 12.6, 0.42, "151B2D")
stats_txt = "  16,831 Total      16,750 ORD      81 Patent      5,017 w/Yield      16,770 High Confidence  "
txt(s, stats_txt, 0.45, 1.69, 12.4, 0.38, size=10, color="94A3B8", font="Calibri")

# Sidebar
box(s, 0.35, 2.09, 2.25, 5.16, "141927", "1E2744", 1)
txt_lines(s, [
    ("SOURCE", True, "64748B", 8),
    ("● All sources", False, WHITE, 10),
    ("○ ORD (expert)", False, "34D399", 10),
    ("○ Patent (LLM)", False, "A78BFA", 10),
    ("", False, WHITE, 6),
    ("YIELD (%)", True, "64748B", 8),
    ("Min: ___  Max: ___", False, "94A3B8", 10),
    ("", False, WHITE, 6),
    ("MIN CONFIDENCE", True, "64748B", 8),
    ("▬▬▬▬▬○────── 0.0", False, "6366F1", 10),
    ("", False, WHITE, 6),
    ("SOLVENT", True, "64748B", 8),
    ("e.g. water, ethanol", False, "475569", 10),
    ("", False, WHITE, 6),
    ("CATALYST", True, "64748B", 8),
    ("e.g. palladium", False, "475569", 10),
    ("", False, WHITE, 6),
    ("SORT BY", True, "64748B", 8),
    ("Confidence ↓", False, "94A3B8", 10),
], 0.45, 2.15, 2.05, 4.8, font="Calibri", line_spacing=15, margin=0.05)

# Main area
box(s, 2.6, 2.09, 10.35, 5.16, "0F1117")

# Search bar
box(s, 2.75, 2.18, 10.05, 0.48, "151B2D", "2A3050", 1)
txt(s, "🔍   Search by keyword, SMARTS, product SMILES, solvent, catalyst...",
    2.85, 2.22, 9.0, 0.4, size=11, color="475569", font="Calibri")
box(s, 11.7, 2.23, 0.85, 0.38, "4F46E5")
txt(s, "Search", 11.7, 2.23, 0.85, 0.38, size=10, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Calibri")

# Result cards
card_data = [
    ("ORD · Expert", "34D399", "0D3B2E", "Sonogashira Coupling",
     "SMARTS: [c:1]-[I]+[C:2]#[CH]>>[c:1]-[C:2]#[CH]",
     "Yield: 99.0%  ·  Solvent: water  ·  Catalyst: Fe3O4@SiO2-NH2-Pd(ii), CuI"),
    ("Patent · LLM  conf=0.60", "A78BFA", "1E1A3B", "Suzuki Cross-Coupling",
     "SMARTS: [c:1]-[Br].[B:2](O)O>>[c:1]-[B:2]",
     "Yield: 87.0%  ·  Temp: 80°C  ·  Solvent: DMF/H2O  ·  Catalyst: Pd(PPh3)4"),
]
cy = 2.82
for (badge, badge_c, badge_bg, title, smarts, meta) in card_data:
    box(s, 2.75, cy, 10.05, 0.95, "1A1F2E", "2A3050", 1)
    box(s, 2.85, cy+0.08, 2.0, 0.28, badge_bg)
    txt(s, badge, 2.87, cy+0.09, 1.96, 0.26, size=8.5, bold=True, color=badge_c, font="Calibri")
    txt(s, title, 5.0, cy+0.06, 4, 0.3, size=11, bold=True, color=WHITE, font="Calibri")
    txt(s, smarts, 2.85, cy+0.42, 10.0, 0.25, size=9, color="6366F1", font="Courier New")
    txt(s, meta, 2.85, cy+0.67, 10.0, 0.22, size=9.5, color="94A3B8", font="Calibri")
    cy += 1.08

txt(s, "Click any SMILES/SMARTS chip to copy  ·  Paginated · 20 results/page",
    2.75, 5.05, 10.0, 0.3, size=9.5, color="475569", font="Calibri", align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Agent Integration
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, LIGHT_BG)

box(s, 0, 0, W, 1.05, NAVY)
txt(s, "ORD Tool — Wired into SynAgent's Agent Network", 0.5, 0.15, 12.3, 0.75,
    size=30, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.LEFT)

# Agent tree — left side
box(s, 0.4, 1.2, 6.0, 5.6, WHITE, "E2E8F0", 1)
txt(s, "Agent Architecture", 0.6, 1.3, 5.6, 0.4, size=13, bold=True, color=DARK, font="Calibri")

# Master
box(s, 0.7, 1.85, 5.4, 0.55, NAVY)
txt(s, "🧠  Master Agent (Gemini 2.5 Flash)", 0.8, 1.87, 5.2, 0.51,
    size=12, bold=True, color=WHITE, font="Calibri", valign=MSO_ANCHOR.MIDDLE)

agents = [
    ("Validation Agent", "RDKit SMILES & template check", "334155"),
    ("ChemSpace Agent", "Building block pricing", "334155"),
    ("Optimization Agent", "PubChem GHS hazard scoring", "334155"),
    ("ORD Agent  ← NEW", "Reaction database lookup (FTS5)", "0A5C5C"),
]
ay = 2.6
for (name, desc, fill) in agents:
    # Connector line
    ln = s.shapes.add_shape(1, Inches(1.15), Inches(ay), Inches(0.04), Inches(0.35))
    ln.fill.solid(); ln.fill.fore_color.rgb = rgb("CBD5E1")
    ln.line.fill.background()
    box(s, 1.25, ay+0.3, 4.7, 0.55, fill)
    txt(s, name, 1.38, ay+0.32, 2.4, 0.51, size=11, bold=True, color=WHITE, font="Calibri",
        valign=MSO_ANCHOR.MIDDLE)
    txt(s, desc, 3.85, ay+0.32, 2.1, 0.51, size=10, color="94A3B8", font="Calibri",
        valign=MSO_ANCHOR.MIDDLE)
    ay += 0.95

# Right side — ORD query flow
box(s, 6.8, 1.2, 6.15, 2.7, "E8F5F3", TEAL, 1)
txt(s, "Example ORD Agent Query", 7.0, 1.3, 5.7, 0.4,
    size=13, bold=True, color=TEAL_BOX, font="Calibri")
txt_lines(s, [
    ('query_ord_reaction(', False, "6366F1", 11),
    ('    reaction_smarts=', False, "94A3B8", 11),
    ('        "[c:1]-[Br].[B:2]>>[c:1]-[B:2]",', False, "34D399", 11),
    ('    keyword="Suzuki palladium"', False, "34D399", 11),
    (')', False, "6366F1", 11),
    ('', False, WHITE, 6),
    ('→ Found 3 reactions', True, "0D9488", 12),
    ('→ 99% yield · water · Pd catalyst', True, "0D9488", 12),
], 7.0, 1.82, 5.7, 2.0, font="Courier New", line_spacing=17, margin=0.06)

box(s, 6.8, 4.1, 6.15, 2.65, WHITE, "E2E8F0", 1)
txt(s, "How It Works", 7.0, 4.2, 5.7, 0.38,
    size=13, bold=True, color=DARK, font="Calibri")
txt_lines(s, [
    ("User: Find precedent for Suzuki coupling", False, MID, 11),
    ("", False, WHITE, 4),
    ("1. Master calls call_reaction_agent()", False, DARK, 11),
    ("2. ORD agent runs FTS5 query on reactions", False, DARK, 11),
    ("3. Returns yield, temp, solvent, catalyst", False, DARK, 11),
    ("4. Master combines with all 4 agent outputs", False, DARK, 11),
    ("5. Single unified synthesis route report", True, TEAL, 12),
], 7.0, 4.65, 5.7, 2.0, font="Calibri", line_spacing=18, margin=0.06)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Engineering Decisions
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, LIGHT_BG)

box(s, 0, 0, W, 1.05, NAVY)
txt(s, "Engineering Decisions", 0.5, 0.15, 12.3, 0.75,
    size=30, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.LEFT)

# Left — decisions
box(s, 0.4, 1.2, 6.0, 6.0, WHITE, "E2E8F0", 1)
txt(s, "What I Chose & Why", 0.6, 1.3, 5.6, 0.4,
    size=14, bold=True, color=DARK, font="Calibri")
decisions = [
    ("SQLite not Postgres", "Portable, runs on USB drives, WAL concurrency mode"),
    ("source column design", "One table, two origins — no complex joins needed"),
    ("FTS5 over vector search", "Fast, no GPU, works fully offline"),
    ("No LLM for ORD data", "Eliminates hallucination, saves API cost"),
    ("run_loop.py supervisor", "Auto-restarts worker on 429 rate-limit failures"),
    ("Confidence scoring", "Flags low-quality extractions for human review"),
    ("PATENT_DATA_DIR env var", "USB path set once; all code reads from it"),
]
dy = 1.85
for (title, desc) in decisions:
    circ = s.shapes.add_shape(9, Inches(0.55), Inches(dy+0.06), Inches(0.28), Inches(0.28))
    circ.fill.solid(); circ.fill.fore_color.rgb = rgb(TEAL)
    circ.line.fill.background()
    txt(s, title, 0.95, dy, 5.2, 0.32, size=12, bold=True, color=DARK, font="Calibri")
    txt(s, desc, 0.95, dy+0.3, 5.2, 0.32, size=11, color=MID, font="Calibri")
    dy += 0.73

# Right — next steps
box(s, 6.9, 1.2, 6.0, 6.0, "F0FDF4", "059669", 1)
txt(s, "Suggested Next Steps", 7.1, 1.3, 5.6, 0.4,
    size=14, bold=True, color="065F46", font="Calibri")
nexts = [
    ("Build FAISS semantic index", "Embedding-based SMILES similarity search"),
    ("Active learning fine-tuning", "Use human corrections to improve LLM extraction"),
    ("PatentsView bulk integration", "Free weekly USPTO grant CSVs (no API key)"),
    ("Deploy search UI to web", "Share with lab — no local setup needed"),
    ("Expand paper corpus", "1,700+ papers already queued for parsing"),
    ("Reaction pathway scoring", "Use ORD precedent to weight SynAgent route scores"),
]
ny = 1.85
for (title, desc) in nexts:
    box(s, 7.05, ny+0.02, 0.28, 0.28, "059669")
    txt(s, "→", 7.05, ny, 0.28, 0.35, size=13, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Calibri")
    txt(s, title, 7.45, ny, 5.2, 0.32, size=12, bold=True, color="065F46", font="Calibri")
    txt(s, desc, 7.45, ny+0.3, 5.2, 0.32, size=11, color="166534", font="Calibri")
    ny += 0.88

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Summary / Closing
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s, NAVY)

# Subtle background shapes
circ = s.shapes.add_shape(9, Inches(-1), Inches(-1), Inches(5), Inches(5))
circ.fill.solid(); circ.fill.fore_color.rgb = rgb("0F2A50")
circ.line.fill.background()
circ2 = s.shapes.add_shape(9, Inches(10), Inches(4.5), Inches(4), Inches(4))
circ2.fill.solid(); circ2.fill.fore_color.rgb = rgb("0A2040")
circ2.line.fill.background()

# Teal left bar
box(s, 0, 0, 0.18, H, TEAL)

# Title
txt(s, "What This Enables", 0.5, 0.8, 12.3, 0.85,
    size=36, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.CENTER)

# Main statement box
box(s, 1.0, 1.85, 11.3, 1.3, "0D1F3C", TEAL, 2)
txt(s, "SynAgent can now answer:\n"
    "Has this reaction been done before? What yield was achieved? Under what conditions?",
    1.2, 1.95, 10.9, 1.1,
    size=17, bold=False, color="7DD3FC", font="Georgia",
    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# Three outcome rows
outcomes = [
    ("✓", "16,831 reactions indexed from ORD + patent literature", "34D399"),
    ("✓", "Any route checked for experimental precedent in milliseconds", "34D399"),
    ("✓", "Human-in-the-loop review improves extraction quality over time", "34D399"),
]
oy = 3.4
for (mark, text, color) in outcomes:
    box(s, 1.2, oy, 0.55, 0.52, "0A3320")
    txt(s, mark, 1.2, oy, 0.55, 0.52, size=18, bold=True, color=color,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, font="Georgia")
    txt(s, text, 1.9, oy+0.07, 10.5, 0.4, size=15, color=WHITE, font="Calibri")
    oy += 0.72

# Divider
box(s, 0.5, 5.35, 12.3, 0.02, TEAL)

# Bottom
txt(s, "github.com/A-Sanil/SynAgent  ·  asanil@berkeley.edu",
    0.5, 5.55, 8.5, 0.45, size=12, color="64748B", font="Calibri")
txt(s, "Excited for Wednesday!",
    8.5, 5.55, 4.5, 0.45, size=13, bold=True, color=TEAL, font="Georgia",
    align=PP_ALIGN.RIGHT)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\Aadit\OneDrive\Documents\GitHub\SynAgent\SynAgent_Aaditya_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
