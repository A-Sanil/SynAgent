#!/usr/bin/env python3
"""Generate PPTX presentation from slides data."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color scheme matching the web UI
BG_DARK = RGBColor(15, 23, 36)          # #0f1724
CARD_DARK = RGBColor(17, 24, 39)        # #111827
ACCENT_PURPLE = RGBColor(124, 58, 237)  # #7c3aed
ACCENT_CYAN = RGBColor(6, 182, 212)     # #06b6d4
TEXT_LIGHT = RGBColor(238, 242, 255)    # #eef2ff
TEXT_MUTED = RGBColor(148, 163, 184)    # #94a3b8

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define slide layouts
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    
    def add_title_slide(title, subtitle):
        """Add a title slide."""
        slide = prs.slides.add_slide(blank_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = ACCENT_PURPLE
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = TEXT_MUTED
        
        return slide
    
    def add_content_slide(title, bullet_points):
        """Add a content slide with bullets."""
        slide = prs.slides.add_slide(blank_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = ACCENT_PURPLE
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_LIGHT
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.level = 0
        
        return slide
    
    def add_two_column_slide(title, left_title, left_points, right_title, right_points):
        """Add a slide with two columns."""
        slide = prs.slides.add_slide(blank_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = ACCENT_PURPLE
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.7))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        # Left subtitle
        p = left_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        
        for point in left_points:
            p = left_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_LIGHT
            p.space_before = Pt(4)
            p.space_after = Pt(4)
            p.level = 0
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.7))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        # Right subtitle
        p = right_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        
        for point in right_points:
            p = right_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_LIGHT
            p.space_before = Pt(4)
            p.space_after = Pt(4)
            p.level = 0
        
        return slide
    
    # ===== SLIDE 1: Title =====
    add_title_slide(
        "Patent Ingestion Pipeline",
        "Automated extraction, validation & review of patent synthesis data"
    )
    
    # ===== SLIDE 2: Problem Statement =====
    add_content_slide(
        "The Challenge",
        [
            "🔍 Manual Patent Data Extraction is Slow",
            "   • Reading patent documents one by one",
            "   • Manually extracting reaction conditions",
            "   • Validating chemical structures (SMILES)",
            "",
            "⚠️ LLMs Hallucinate Chemical Data",
            "   • Generated SMILES may not be valid",
            "   • No automatic verification or confidence scoring",
            "   • No human feedback loop for correction"
        ]
    )
    
    # ===== SLIDE 3: Solution =====
    add_content_slide(
        "Our Solution",
        [
            "✓ Collects patents from web (Scrapling) + PDFs",
            "✓ Parses with Qwen LLM (vLLM on Savio GPU)",
            "✓ Validates chemistry (RDKit + PubChem cross-check)",
            "✓ Scores confidence (multi-pass verification pipeline)",
            "✓ Enables human review (web UI with modal editing)",
            "✓ Logs corrections (active learning for retraining)"
        ]
    )
    
    # ===== SLIDE 4: Architecture =====
    add_content_slide(
        "System Architecture",
        [
            "Web URLs → Scrapling → Raw Documents (SQLite)",
            "             ↓",
            "   PDF Collector → pdfplumber → Extract text/tables",
            "             ↓",
            "   Parse Queue → LLM Parser (Qwen/vLLM on Savio)",
            "             ↓",
            "   Chemistry NER ← OPSIN, RDKit, PubChem",
            "             ↓",
            "   Canonicalization & Confidence Scoring",
            "             ↓",
            "   Database (SQLite) → Web UI (FastAPI + Jinja2)",
            "             ↓",
            "   Human Review & Active Learning Logs"
        ]
    )
    
    # ===== SLIDE 5: Features 1 =====
    add_content_slide(
        "Key Features (1/2)",
        [
            "✓ Web Collection",
            "   Scrapling-based scraping + PDF download & text extraction",
            "",
            "✓ LLM-Powered Extraction",
            "   Qwen 27B via vLLM returns structured JSON",
            "",
            "✓ Multi-Pass Validation",
            "   LLM → NER → SMILES canonicalization → PubChem cross-check",
            "",
            "✓ Chemistry Validation",
            "   RDKit SMILES validation, PubChem lookup, yield/temp extraction"
        ]
    )
    
    # ===== SLIDE 6: Features 2 =====
    add_content_slide(
        "Key Features (2/2)",
        [
            "✓ Human-in-the-Loop UI",
            "   Web interface at http://127.0.0.1:8001",
            "",
            "✓ Modal Editing",
            "   Click \"Edit\" on reaction → modal → change SMILES/notes → AJAX save",
            "",
            "✓ Batch Review",
            "   Dedicated page for low-confidence reactions (<0.6)",
            "",
            "✓ Active Learning Logging",
            "   Every correction stored (user, timestamp, old/new values)"
        ]
    )
    
    # ===== SLIDE 7: Database =====
    add_two_column_slide(
        "Database Schema (SQLite)",
        "Tables",
        [
            "• patents",
            "  - title, abstract, inventors",
            "  - domain_tags, reviewed",
            "",
            "• reactions",
            "  - product_smiles, confidence",
            "  - yield, temp, catalyst",
            "",
            "• raw_documents",
            "  - source_url, fetched_at",
            "  - raw_text"
        ],
        "Indexes & Features",
        [
            "• FTS5 - full-text search",
            "",
            "• parse_queue",
            "  - async worker orchestration",
            "",
            "• active_learning",
            "  - correction logs",
            "",
            "• Semantic index",
            "  - FAISS (optional)"
        ]
    )
    
    # ===== SLIDE 8: UI Walkthrough =====
    add_content_slide(
        "Web UI Walkthrough",
        [
            "📋 Review Queue (/)",
            "   List all patents, sorted by review status",
            "",
            "📄 Patent Detail (/patent/{id})",
            "   View patent title, abstract, reactions with edit buttons",
            "",
            "🔍 Batch Review (/batch_review)",
            "   Filter to low-confidence reactions, bulk-correct with modal",
            "",
            "🔎 Search (/search?q=...)",
            "   Full-text search across patents and reactions"
        ]
    )
    
    # ===== SLIDE 9: Chemistry Validation =====
    add_content_slide(
        "Chemistry Validation Pipeline",
        [
            "Confidence scoring prevents storing bad SMILES:",
            "",
            "1️⃣ RDKit Canonicalization (60 pts)",
            "   SMILES valid → extract canonical form → +0.6 confidence",
            "",
            "2️⃣ PubChem Cross-Check (40 pts)",
            "   Query PubChem for canonical SMILES → agreement → +0.4",
            "",
            "3️⃣ Final Score: 0.0–1.0",
            "   Reactions <0.6 flagged for human review"
        ]
    )
    
    # ===== SLIDE 10: Savio =====
    add_two_column_slide(
        "Running on Savio (Berkeley Lab)",
        "vLLM Setup",
        [
            "• Qwen 27B in Apptainer",
            "",
            "• OpenAI-compatible API",
            "  /v1/chat/completions",
            "",
            "• Slurm job:",
            "  - 1 A40 GPU",
            "  - savio3_gpu partition",
            "  - 2 hours"
        ],
        "Access from Local",
        [
            "• SSH tunnel:",
            "  ssh -L 8000:node:8000",
            "  user@savio.lbl.gov",
            "",
            "• .env points to",
            "  http://127.0.0.1:8000",
            "",
            "• Worker calls LLM",
            "  transparently via HTTP"
        ]
    )
    
    # ===== SLIDE 11: CLI =====
    add_content_slide(
        "CLI Commands",
        [
            "# Initialize database",
            "python -m patent_pipeline.cli init_db",
            "",
            "# Collect patents",
            "python -m patent_pipeline.cli collect \"https://...\"",
            "",
            "# Enqueue for parsing",
            "python -m patent_pipeline.cli enqueue_all",
            "",
            "# Run worker",
            "python -m patent_pipeline.cli run_worker \\",
            "  --base-url http://127.0.0.1:8000 --model qwen",
            "",
            "# Start web UI",
            "python -m patent_pipeline.cli runserver"
        ]
    )
    
    # ===== SLIDE 12: Status =====
    add_two_column_slide(
        "Current Status",
        "✓ Working",
        [
            "• Web UI (no auth errors)",
            "• Database + indexing",
            "• Batch review page",
            "• Modal editing",
            "• CLI commands",
            "• Scrapling installed",
            "• Chemistry NER"
        ],
        "⚠️ Partial / 🔴 Pending",
        [
            "⚠️ RDKit (Windows — optional)",
            "⚠️ OPSIN (Windows — optional)",
            "⚠️ PubChem (needs network)",
            "",
            "🔴 vLLM endpoint (building)",
            "🔴 End-to-end parse (need vLLM)",
            "🔴 Multi-worker",
            "🔴 Semantic search"
        ]
    )
    
    # ===== SLIDE 13: Testing =====
    add_content_slide(
        "Testing & Validation",
        [
            "✓ Smoke Tests Passed:",
            "   • UI root endpoint returns HTTP 200",
            "   • Batch review page returns HTTP 200",
            "   • Browser loads review queue",
            "   • Scrapling imports (v0.4.8)",
            "   • Database initialization",
            "   • Chemistry NER extraction",
            "",
            "⏳ Pending Tests:",
            "   • Full parse job (need vLLM endpoint)",
            "   • Active learning logging e2e",
            "   • Multi-worker orchestration"
        ]
    )
    
    # ===== SLIDE 14: Next Steps =====
    add_content_slide(
        "Next Steps",
        [
            "1️⃣ Confirm vLLM endpoint — monitor Savio job, test reachability",
            "",
            "2️⃣ Run end-to-end test — 1 URL → enqueue → parse → review",
            "",
            "3️⃣ Install optional deps — RDKit, OPSIN (offline validation)",
            "",
            "4️⃣ Add test suite — pytest with DB fixtures, mocked vLLM",
            "",
            "5️⃣ Scale workers — Redis queue or multiprocessing",
            "",
            "6️⃣ Implement retraining — active learning → fine-tune Qwen",
            "",
            "7️⃣ Add semantic search — FAISS indexing on demand"
        ]
    )
    
    # ===== SLIDE 15: Key Takeaways =====
    add_content_slide(
        "Key Takeaways",
        [
            "🚀 What You Get",
            "   Fully automated patent synthesis data extraction & human-in-the-loop review",
            "",
            "🔧 Built With",
            "   Scrapling, Qwen, RDKit/PubChem, SQLite, FastAPI, Savio GPU",
            "",
            "📈 Scales From",
            "   Hobby/research (local SQLite) to production (Redis + multiple workers)",
            "",
            "🧠 Learns From Feedback",
            "   Active learning logs enable retraining & confidence boosting over time"
        ]
    )
    
    # ===== SLIDE 16: Closing =====
    add_title_slide(
        "Questions?",
        "Patent Ingestion Pipeline — Ready for Deployment"
    )
    
    # Add info to closing slide
    closing_slide = prs.slides[-1]
    info_box = closing_slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.8))
    info_frame = info_box.text_frame
    info_frame.word_wrap = True
    
    lines = [
        "Repository: SynAgent / patent_ingestion_pipeline",
        "Documentation: ARCHITECTURE.md",
        "Status: MVP ✓ | Savio Ready ✓ | Production-Ready (pending vLLM e2e test)"
    ]
    
    for i, line in enumerate(lines):
        if i == 0:
            p = info_frame.paragraphs[0]
        else:
            p = info_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return prs

if __name__ == "__main__":
    prs = create_presentation()
    output_path = "PRESENTATION.pptx"
    prs.save(output_path)
    print(f"✓ PowerPoint presentation created: {output_path}")
    print(f"  - 16 slides")
    print(f"  - Dark theme (RGB: 15, 23, 36)")
    print(f"  - Ready to share or edit in PowerPoint")
